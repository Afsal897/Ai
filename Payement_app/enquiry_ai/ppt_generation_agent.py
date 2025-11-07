import os
import json
import logging
from datetime import datetime, timezone
from typing import List, Dict, Optional
import uuid
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from config import Config
from pydantic import BaseModel, Field
from langchain_core.tools import tool
from langchain_google_genai import ChatGoogleGenerativeAI
from sqlalchemy.engine.base import Engine
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from log import logger
# ---------------------------
# Logging
# ---------------------------


KEY_FEATURES = "Key Features"
THANK_YOU = "Thank You"

# ---------------------------
# Gemini LLM Wrapper
# ---------------------------
class GeminiLLM:
    def __init__(self):
        if not Config.GOOGLE_API_KEY_1:
            raise ValueError("GOOGLE_API_KEY_1 is required")
        self.model = ChatGoogleGenerativeAI(
            model="gemini-2.5-flash",
            google_api_key=Config.GOOGLE_API_KEY_1,
            temperature=0.2,
        )

    def create_slide_outline(
        self,
        user_instruction: str,
        retrieved_texts: List[Dict],
    ) -> List[Dict]:

        context_block = retrieved_texts

        # ✅ COMBINED SYSTEM PROMPT
        system_prompt = """
         You are a specialized AI assistant that transforms unstructured project data into structured JSON presentation outlines.

        Your primary mission is to analyze the provided project data, identify distinct project groups based on shared domains or technologies, and generate a unique 5-slide presentation outline for each identified group.

        You must follow these rules with absolute precision.

        ### 1. Strict JSON Output Format
        - The entire response **MUST** be a single, valid JSON object. Do not add any text or markdown formatting outside this JSON structure.
        - The root object must contain exactly two keys: suggested_filename and slides.
        - **suggested_filename**: A string containing the generated filename with a timestamp or UUID.
        - **slides**: A flat list of all slide objects from every generated presentation.

        ### 2. Confidentiality Mandate
        - **CRITICAL**: Under no circumstances should any specific project names or client names appear in the output.
        - All bullet points must contain generalized, anonymized information synthesized from the source data.

        ### 3. Presentation Outline Structure (Per Group)
        - Each project group must be represented by an outline of **EXACTLY 5 slides**.
        - The title for the 5 slides must follow this precise order:
            1. **Title** — Use the project name from the user instruction (or derived group title if multiple projects).
            2. Overview
            3. Key Features
            4. Technologies
            5. Thank You

        ### 4. Slide Content Rules
        - **Title & Thank You Slides (Slides 1 & 5):** The bullets array for these slides must be empty.
        - **Content Slides (Slides 2, 3, & 4):** Each must contain between 3–6 concise bullet points (point-wise, not paragraphs).
        - The content must be a synthesized, clear, and professional summary aggregated from relevant data.

        ### 5. Example Output Structure
        {
          "suggested_filename": "project_presentation_20251009_153045.json",
          "slides": [
            {"title": "Project Name", "bullets": []},
            {"title": "Overview", "bullets": ["Point 1", "Point 2", "Point 3"]},
            {"title": "Key Features", "bullets": ["Point 1", "Point 2", "Point 3"]},
            {"title": "Technologies", "bullets": ["Point 1", "Point 2", "Point 3"]},
            {"title": "Thank You", "bullets": []}
          ]
        }
        
        """

        human_prompt = (
            f"User instruction:\n{user_instruction}\n\n"
            f"Context excerpts:\n{context_block}\n\n"
            "Constraints:\n"
            "- Aggregate all projects under the same domain, technology, or sub-technology into a single PPT per group.\n"
            "- Maximum slides per PPT: 5.\n"
            "- Slide titles must be exactly: Title, Overview, Key Features, Technologies, Thank You.\n"
            "- The Title slide must use the project name provided in the user instruction.\n"
            "- Slides 2–4 must include 3–6 concise, point-wise bullet points (not paragraphs).\n"
            "- Slides 1 and 5 must have no bullets.\n"
            "- Maintain confidentiality: do not include specific client or project names.\n"
            "- Respond strictly in valid JSON format only, with the specified slide structure. Do not add explanations or extra text.\n"
            "- Include a 'suggested_filename' key with a unique timestamp or UUID."
        )

        response = self.model.invoke(system_prompt + "\n\n" + human_prompt)
        text_out = getattr(response, "content", str(response)).strip()
        print(f"response--->{response}")

        try:
            logger.info("Extracting JSON from response.")
            first_brace = text_out.find("{")
            last_brace = text_out.rfind("}") + 1
            json_text = text_out[first_brace:last_brace]
            parsed = json.loads(json_text)

            slides = parsed.get("slides", [])
            clean_slides = []
            for s in slides[:5]:
                title = (s.get("title") or "").strip()
                bullets = [str(b).strip() for b in s.get("bullets", [])][:6]
                if title:
                    clean_slides.append({"title": title, "bullets": bullets})

            if len(clean_slides) < 5:
                default_titles = ["Title", "Overview", KEY_FEATURES, "Technologies", THANK_YOU]
                for t in default_titles[len(clean_slides):]:
                    clean_slides.append({"title": t, "bullets": []})

            return clean_slides

        except Exception as e:
            logger.warning("LLM response invalid, using fallback: %s", e)
            return [
                {"title": "Title", "bullets": []},
                {"title": "Overview", "bullets": ["High-level summary of the project."]},
                {"title": KEY_FEATURES, "bullets": ["Feature 1", "Feature 2", "Feature 3"]},
                {"title": "Technologies", "bullets": ["Python", "Flask", "SQLAlchemy"]},
                {"title": THANK_YOU, "bullets": []},
            ]


# ---------------------------
# PPT Generator
# ---------------------------
class PPTGenerator:
    def __init__(self, template_path: str, output_folder: str):
        self.template_path = template_path
        self.output_folder = output_folder
        os.makedirs(self.output_folder, exist_ok=True)

    def _remove_all_slides(self, prs: Presentation):
        xml_slides = prs.slides._sldIdLst
        slide_ids = list(xml_slides)
        for slide_id in slide_ids:
            rel_id = slide_id.rId
            prs.part.drop_rel(rel_id)
            xml_slides.remove(slide_id)

    def _load_presentation(self) -> Presentation:
        if os.path.exists(self.template_path):
            prs = Presentation(self.template_path)
            if prs.slides:
                self._remove_all_slides(prs)
            return prs
        return Presentation()

    def _find_body_placeholder(self, slide):
        for ph in slide.placeholders:
            try:
                if ph.placeholder_format.type == PP_PLACEHOLDER.BODY and ph.has_text_frame:
                    return ph
            except Exception:
                continue
        for ph in slide.placeholders:
            if ph.has_text_frame and ph != slide.shapes.title:
                return ph
        return None

    def _add_bullets(self, text_frame, bullets: list):
        for b in bullets:
            if not b.strip():
                continue
            p = text_frame.add_paragraph()
            p.text = f"• {b.strip()}"
            p.level = 0
            p.font.size = Pt(18)
            p.space_after = Pt(5)

    def _add_slide(self, prs: Presentation, slide_data: dict, is_title_slide=False):
        slide_title = slide_data.get("title", "").strip()
        title_font_size = Pt(44)

        # ---------------- Title Slide ----------------
        if is_title_slide:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2),
                prs.slide_width - Inches(2), Inches(2)
            )
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = slide_title or "Presentation Title"
            p.font.size = title_font_size
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            return

        # ---------------- Thank You Slide ----------------
        if slide_title.lower() == "thank you":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            # Remove all placeholders
            for shape in list(slide.shapes):
                if hasattr(shape, "text_frame"):
                    slide.shapes._spTree.remove(shape._element)

            # Add centered "Thank You" text
            text_box = slide.shapes.add_textbox(
                Inches(1), prs.slide_height / 2 - Inches(1),
                prs.slide_width - Inches(2), Inches(2)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = title_font_size
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            return

        # ---------------- Normal Content Slides ----------------
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
            p = slide.shapes.title.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = title_font_size
            p.font.bold = True

        # Remove all non-title placeholders
        for shape in list(slide.shapes):
            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                slide.shapes._spTree.remove(shape._element)
        # Common constants for all content slides
        LEFT_MARGIN = Inches(1)
        TOP_MARGIN = Inches(1.2)   # distance from top of slide
        BOX_WIDTH = prs.slide_width - Inches(2)  # leave 1 inch margin both sides
        BOX_HEIGHT = prs.slide_height - TOP_MARGIN - Inches(1)  # leave 1 inch bottom margin

        # Add top-left aligned textbox for bullets
        text_box = slide.shapes.add_textbox(
            LEFT_MARGIN, TOP_MARGIN,
            BOX_WIDTH, BOX_HEIGHT
        )
        tf = text_box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True

        # Add bullets (truncate if too long)
        MAX_CHARS = 120
        for b in slide_data.get("bullets", []):
            if not b.strip():
                continue
            text = b.strip()
            if len(text) > MAX_CHARS:
                text = text[:MAX_CHARS-3] + "..."
            p = tf.add_paragraph()
            p.text = f"• {text}"
            p.level = 0
            p.font.size = Pt(18)
            p.space_after = Pt(5)
            p.alignment = PP_ALIGN.LEFT


      
        


    def create_ppt(self, slides: list, output_name: Optional[str] = None) -> str:
        prs = self._load_presentation()
        for idx, slide_data in enumerate(slides):
            is_title = idx == 0  # first slide is title
            self._add_slide(prs, slide_data, is_title_slide=is_title)

        timestamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        uid = uuid.uuid4().hex[:6]
        out_name = (output_name or "presentation.pptx").replace(".pptx", f"_{timestamp}_{uid}.pptx")
        out_path = os.path.join(self.output_folder, out_name)
        prs.save(out_path)
        logger.info(f"PPT saved at: {out_path}")
        return out_path


# ---------------------------
# PPT Service
# ---------------------------
class PPTService:
    def __init__(self, db_engine: Optional[Engine] = None):
        self.llm = GeminiLLM()
        self.pptgen = PPTGenerator(Config.TEMPLATE_PATH, Config.OUTPUT_FOLDER)

    def generate_ppt(
        self,
        instruction: str,
        context: Optional[str] = None,
        slides_json: Optional[dict] = None,
        output_name: Optional[str] = None,
    ) -> str:
        if slides_json:
            slides = slides_json.get("slides", [])
        else:
            retrieved = [{"text": context}] if context else []
            slides = self.llm.create_slide_outline(instruction, retrieved)
        return self.pptgen.create_ppt(slides, output_name=output_name)


# ---------------------------
# LangChain Tool
# ---------------------------
ppt_service: Optional[PPTService] = None

def init_ppt_service(db_engine: Optional[Engine] = None):
    global ppt_service
    ppt_service = PPTService(db_engine)


class GeneratePPTArgs(BaseModel):
    instruction: str = Field(..., description="High-level instruction for the PPT (e.g., what to summarize)")
    context: Optional[str] = Field(None, description="Raw text or JSON string from SQL agent to use as context")
    slides_json: Optional[str] = Field(None, description="Strict JSON with {'slides': [...]} to bypass LLM")
    output_name: Optional[str] = Field(None, description="Optional output filename, e.g., 'my_slides.pptx'")


@tool("generate_ppt_tool", args_schema=GeneratePPTArgs)
def generate_ppt_tool(
    instruction: str,
    context: Optional[str] = None,
    slides_json: Optional[str] = None,
    output_name: Optional[str] = None
) -> str:
    """
    Generate a PowerPoint (.pptx) presentation from an instruction and optional context or slide JSON.
    First slide title is centered vertically and horizontally.
    """
    if ppt_service is None:
        raise RuntimeError("PPTService is not initialized. Call init_ppt_service(engine) first.")

    parsed_slides = json.loads(slides_json) if slides_json else None
    logger.info("Generating PPT with instruction: %s", instruction) 
    return ppt_service.generate_ppt(
        instruction=instruction,
        context=context,
        slides_json=parsed_slides,
        output_name=output_name,
    )
