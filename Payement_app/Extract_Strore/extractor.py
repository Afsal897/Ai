from typing import Any
from pptx import Presentation

from log import logger
def extract_text_with_fontsizes(shape):
    """Extract (text, font_size) pairs from a shape."""
    results = []
    logger.debug(f"Extracting text from shape of type {shape.shape_type}")
    if not getattr(shape, "has_text_frame", False):
        return results
    
    try:

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                txt = run.text.strip()
                if not txt:
                    continue
                font_size = None
                if run.font.size:
                    font_size = run.font.size.pt
                results.append((txt, font_size))

    except Exception as e:
        logger.error(f"Error extracting text from shape: {e}")
    return results


def extract_table(shape):
    """Extract table content as list of rows (list of cell texts)."""
    logger.debug("Extracting table from shape.")
    table_data = []
    try:

        tbl = shape.table
        for row in tbl.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)

    except Exception as e:
        logger.warning(f"Skipping shape table due to error: {e}")
    return table_data


def classify_text(fs: int | None) -> str:
    """Classify text based on font size."""
    if not fs or fs <= 24:
        return "body"
    if 25 <= fs <= 39:
        return "subtitle"
    return "title"


def parse_slide(slide) -> dict[str, Any]:
    """Parse one slide into structured dict using font-size thresholds and tables."""
    logger.debug("Parsing slide.")
    categorized_text = {"body": [], "subtitle": [], "title": []}
    tables = []

    for shp in slide.shapes:
        # Text content
        try:

            if shp.has_text_frame:
                for txt, fs in extract_text_with_fontsizes(shp):
                    category = classify_text(fs)
                    categorized_text[category].append(txt)

            # Table content
            if shp.shape_type == 19:
                table_data = extract_table(shp)
                if table_data:
                    tables.append(table_data)

        except Exception as e:
            logger.warning(f"Skipping shape due to unexpected error: {e}")

    return {
        "title": " ".join(categorized_text["title"]).strip(),
        "subtitle": " ".join(categorized_text["subtitle"]).strip(),
        "text_blocks": categorized_text["body"],
        "tables": tables,
    }

def parse_pptx(file_path):
    logger.info(f"Starting PPTX parsing: {file_path}")
    slides_data = []

    try:  

        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        logger.info(f"Total slides found: {total_slides}")
        
        for i in range(total_slides - 1):  # skip last slide
            slide_info = parse_slide(prs.slides[i])
            slides_data.append({
                "slide_number": i + 1,
                **slide_info
            })

        logger.info(f"Finished parsing {len(slides_data)} slides from {file_path}")
    except Exception as e:
        logger.error(f"Failed to parse PPTX {file_path}: {e}")


    return slides_data